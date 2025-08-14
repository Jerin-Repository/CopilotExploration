from pptx import Presentation
import sys
import os


def combine_pptx(pptx_files, output_file):
    """
    Combines multiple PowerPoint (.pptx) files into a single presentation.

    Args:
        pptx_files (list of str): List of file paths to the PPTX files to be combined.
        output_file (str): Path to the output PPTX file where the combined presentation will be saved.

    Returns:
        None

    Notes:
        - The function starts with the first presentation in the list and appends slides from subsequent presentations.
        - If no PPTX files are provided, the function prints a message and returns.
        - The combined presentation is saved to the specified output file.
    """
    if not pptx_files:
        print("No PPTX files provided.")
        return

    # Start with the first presentation
    combined = Presentation(pptx_files[0])

    # For each subsequent file, copy slides into combined
    for pptx_path in pptx_files[1:]:
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            # Copy slide
            slide_xml = slide.element
            new_slide = combined.slides._sldIdLst.addnext(slide_xml)
    combined.save(output_file)
    print(f"Combined presentation saved as {output_file}")

if __name__ == "__main__":
    # Usage: python combine_ppt.py file1.pptx file2.pptx ... output.pptx
    if len(sys.argv) < 4:
        print("Usage: python combine_ppt.py file1.pptx file2.pptx ... output.pptx")
        sys.exit(1)
    pptx_files = sys.argv[1:-1]
    output_file = sys.argv[-1]
    combine_pptx(pptx_files, output_file)    pip install python-pptx